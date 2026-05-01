"""Layer 4 parsing + chunking tests.

Fixtures generate test files on the fly (no binary files in the repo).
OCR test uses RapidOCR; first run downloads ~80 MB ONNX models into the
RapidOCR cache.
"""
from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

import docx as python_docx

from app.services.chunking import build_chunk_pipeline, chunk_documents
from app.services.parsing import parse_file
from app.services.parsing.base import Document
from app.services.parsing.ocr import OCREngine
from app.services.parsing.router import UnsupportedFileType


# ─────────────────────────────────────────────────────────────────────
# Fixtures: build sample files in tmp dir
# ─────────────────────────────────────────────────────────────────────


@pytest.fixture
def text_pdf(tmp_path: Path) -> Path:
    """A text-only PDF (pypdf can extract directly)."""
    p = tmp_path / "text.pdf"
    c = canvas.Canvas(str(p), pagesize=letter)
    c.setFont("Helvetica", 14)
    c.drawString(72, 720, "Annual Leave Policy")
    c.drawString(72, 700, "Employees are entitled to 10 days of paid annual leave per year.")
    c.drawString(72, 680, "Leave accrues monthly and rolls over at fiscal year end.")
    c.showPage()
    c.drawString(72, 720, "Page 2: Sick Leave")
    c.drawString(72, 700, "Up to 5 days per calendar year, doctor's note required after 2.")
    c.save()
    return p


@pytest.fixture
def scanned_pdf(tmp_path: Path) -> Path:
    """A PDF whose content is an image of text — pypdf cannot extract,
    so PdfReader must fall back to OCR."""
    p = tmp_path / "scanned.pdf"

    img = Image.new("RGB", (1000, 240), "white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("arial.ttf", 56)
    except OSError:
        font = ImageFont.load_default()
    draw.text((30, 80), "HELLO WORLD 2026", fill="black", font=font)

    img_path = tmp_path / "_scan.png"
    img.save(img_path)

    c = canvas.Canvas(str(p), pagesize=letter)
    c.drawImage(ImageReader(str(img_path)), 50, 400, width=500, height=120)
    c.save()
    return p


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    p = tmp_path / "sample.docx"
    doc = python_docx.Document()
    doc.add_heading("Product Manual", level=1)
    doc.add_paragraph("This is the introduction paragraph.")
    doc.add_heading("Permissions", level=2)
    doc.add_paragraph("Roles include admin, editor, and viewer.")
    doc.add_heading("Admin", level=3)
    doc.add_paragraph("Admins can do anything.")
    # Table
    tbl = doc.add_table(rows=2, cols=2)
    tbl.rows[0].cells[0].text = "Role"
    tbl.rows[0].cells[1].text = "Can Delete"
    tbl.rows[1].cells[0].text = "admin"
    tbl.rows[1].cells[1].text = "yes"
    doc.save(str(p))
    return p


@pytest.fixture
def sample_md(tmp_path: Path) -> Path:
    p = tmp_path / "guide.md"
    p.write_text(
        "# Architecture Guide\n\n"
        "Welcome to the system overview.\n\n"
        "## Components\n\n"
        "The system has three components: API, worker, and database.\n\n"
        "### API Layer\n\n"
        "FastAPI serves requests over HTTP and SSE.\n\n"
        "### Worker Layer\n\n"
        "Celery executes async ingestion tasks.\n",
        encoding="utf-8",
    )
    return p


@pytest.fixture
def sample_txt(tmp_path: Path) -> Path:
    p = tmp_path / "notes.txt"
    p.write_text("Just some plain notes.\nLine two.\n", encoding="utf-8")
    return p


# ─────────────────────────────────────────────────────────────────────
# Reader / Router
# ─────────────────────────────────────────────────────────────────────


class TestReaders:
    def test_pdf_text_path(self, text_pdf: Path):
        docs = parse_file(text_pdf)
        assert len(docs) == 2  # two pages
        assert "Annual Leave" in docs[0].text
        assert "Sick Leave" in docs[1].text
        for i, d in enumerate(docs, start=1):
            assert d.metadata["parser"] == "pdf"
            assert d.metadata["page"] == i
            assert d.metadata["file_type"] == "pdf"
            assert d.metadata["source"] == "text.pdf"

    def test_pdf_ocr_fallback(self, scanned_pdf: Path):
        # Trigger OCR: should produce a doc tagged "pdf+ocr" with "HELLO" in text
        docs = parse_file(scanned_pdf)
        assert len(docs) >= 1
        d = docs[0]
        assert d.metadata["parser"] == "pdf+ocr"
        # OCR may not be 100% perfect; require at least one of the strong words
        assert any(tok in d.text.upper() for tok in ("HELLO", "WORLD", "2026"))

    def test_docx(self, sample_docx: Path):
        docs = parse_file(sample_docx)
        assert len(docs) == 1
        text = docs[0].text
        # Heading levels reconstructed as markdown
        assert "# Product Manual" in text
        assert "## Permissions" in text
        assert "### Admin" in text
        # Body content present
        assert "introduction paragraph" in text
        # Table rendered as markdown table
        assert "| Role | Can Delete |" in text
        assert "| admin | yes |" in text
        assert docs[0].metadata["parser"] == "docx"

    def test_markdown(self, sample_md: Path):
        docs = parse_file(sample_md)
        assert len(docs) == 1
        assert "# Architecture Guide" in docs[0].text
        assert docs[0].metadata["parser"] == "markdown"

    def test_txt(self, sample_txt: Path):
        docs = parse_file(sample_txt)
        assert len(docs) == 1
        assert "plain notes" in docs[0].text
        assert docs[0].metadata["parser"] == "text"

    def test_html(self, tmp_path: Path):
        p = tmp_path / "page.html"
        p.write_text(
            "<html><head><title>T</title></head><body>"
            "<h1>Hello</h1><p>some content</p>"
            "<h2>Sub</h2><p>more</p>"
            "<script>alert('drop me')</script>"
            "</body></html>",
            encoding="utf-8",
        )
        docs = parse_file(p)
        assert len(docs) == 1
        text = docs[0].text
        assert "# Hello" in text
        assert "## Sub" in text
        assert "some content" in text
        assert "alert" not in text  # script removed
        assert docs[0].metadata["parser"] == "html"

    def test_xlsx(self, tmp_path: Path):
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["类别", "留存周期"])
        ws.append(["用户反馈", "不少于 18 个月"])
        ws.append(["质量抽检", "不少于 2 年"])
        p = tmp_path / "table.xlsx"
        wb.save(str(p))

        docs = parse_file(p)
        assert len(docs) == 1
        text = docs[0].text
        assert "## Sheet1" in text
        assert "| 类别 | 留存周期 |" in text
        assert "| 用户反馈 | 不少于 18 个月 |" in text
        assert docs[0].metadata["parser"] == "xlsx"

    def test_csv(self, tmp_path: Path):
        p = tmp_path / "data.csv"
        p.write_text("name,age\nalice,30\nbob,28\n", encoding="utf-8")
        docs = parse_file(p)
        assert len(docs) == 1
        assert "| name | age |" in docs[0].text
        assert "| alice | 30 |" in docs[0].text

    def test_json(self, tmp_path: Path):
        import json
        p = tmp_path / "config.json"
        p.write_text(json.dumps({"name": "alice", "tags": ["a", "b"]},
                                ensure_ascii=False), encoding="utf-8")
        docs = parse_file(p)
        assert len(docs) == 1
        assert "alice" in docs[0].text
        assert docs[0].metadata["parser"] == "json"

    def test_pptx(self, tmp_path: Path):
        from pptx import Presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Quarterly Review"
        slide.placeholders[1].text = "Revenue grew 12%."
        p = tmp_path / "deck.pptx"
        prs.save(str(p))

        docs = parse_file(p)
        assert len(docs) == 1
        text = docs[0].text
        assert "Quarterly Review" in text
        assert "Revenue grew 12%" in text
        assert docs[0].metadata["parser"] == "pptx"

    def test_image_ocr(self, tmp_path: Path):
        # Reuse the OCR test fixture: render text into a PNG, parse it.
        img = Image.new("RGB", (1000, 200), "white")
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 60)
        except OSError:
            font = ImageFont.load_default()
        draw.text((30, 60), "INVOICE 2026", fill="black", font=font)
        p = tmp_path / "scan.png"
        img.save(p)
        docs = parse_file(p)
        assert len(docs) == 1
        assert any(t in docs[0].text.upper() for t in ("INVOICE", "2026"))
        assert docs[0].metadata["parser"] == "image+ocr"

    def test_unsupported_extension(self, tmp_path: Path):
        bad = tmp_path / "what.xyz"
        bad.write_text("hi")
        with pytest.raises(UnsupportedFileType):
            parse_file(bad)

    def test_missing_file(self, tmp_path: Path):
        with pytest.raises(FileNotFoundError):
            parse_file(tmp_path / "nope.pdf")


# ─────────────────────────────────────────────────────────────────────
# OCR engine direct test
# ─────────────────────────────────────────────────────────────────────


class TestOCREngine:
    def test_recognize_image(self, tmp_path: Path):
        img = Image.new("RGB", (1000, 200), "white")
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 60)
        except OSError:
            font = ImageFont.load_default()
        draw.text((30, 60), "OPEN AI 2025", fill="black", font=font)
        out = OCREngine.get().recognize(img)
        assert isinstance(out, str)
        assert len(out) > 0
        assert any(t in out.upper() for t in ("OPEN", "AI", "2025"))


# ─────────────────────────────────────────────────────────────────────
# Chunking pipeline
# ─────────────────────────────────────────────────────────────────────


class TestChunking:
    def test_pipeline_builds(self):
        p = build_chunk_pipeline()
        # 4 transformations: parser + splitter + table-merger + enricher
        assert len(p.transformations) == 4

    def test_chunks_have_required_metadata(self, sample_md: Path):
        docs = parse_file(sample_md)
        # set doc_id explicitly (Layer 5 will do this)
        for d in docs:
            d.metadata["doc_id"] = "doc-md-1"
        nodes = chunk_documents(docs)
        assert len(nodes) >= 1
        for n in nodes:
            assert "chunk_id" in n.metadata
            assert "chunk_index" in n.metadata
            assert "doc_id" in n.metadata
            assert n.metadata["doc_id"] == "doc-md-1"
            assert "prev_chunk_id" in n.metadata
            assert "next_chunk_id" in n.metadata
            assert "breadcrumbs" in n.metadata

    def test_breadcrumbs_prepended(self, sample_md: Path):
        docs = parse_file(sample_md)
        for d in docs:
            d.metadata["doc_id"] = "doc-md-1"
        nodes = chunk_documents(docs)
        # at least one node has breadcrumbs and starts with "["
        crumbed = [n for n in nodes if n.metadata.get("breadcrumbs")]
        assert crumbed, "expected at least one chunk with breadcrumbs"
        for n in crumbed:
            assert n.text.startswith("["), f"text should start with bread crumb prefix: {n.text[:60]!r}"

    def test_markdown_table_not_split_across_chunks(self, tmp_path: Path):
        """Regression: SentenceSplitter used to cut markdown tables mid-row,
        producing chunks whose first line was the *trailing cell* of the
        previous row — LLMs then aligned that value to the next row's label
        and gave confidently-wrong answers. The MarkdownTableMerger restitches
        the table so every row's key and value live in the same chunk."""
        from app.services.chunking import build_chunk_pipeline

        # Big enough that chunk_size=512 will definitely cut it
        rows = []
        for i in range(40):
            rows.append(f"| 类别 {i:02d} | 一些较长的描述文本，用于撑大每行的 token 数，让分块器在表格中间触发切割 | 不少于 {i + 1} 年 |")
        text = (
            "# 标题\n\n"
            "前置段落，让表头有点上下文。\n\n"
            "| 类别 | 描述 | 留存周期 |\n"
            "| --- | --- | --- |\n"
            + "\n".join(rows)
            + "\n\n后置段落。\n"
        )
        path = tmp_path / "table.md"
        path.write_text(text, encoding="utf-8")

        docs = parse_file(path)
        for d in docs:
            d.metadata["doc_id"] = "doc-table-1"
        pipeline = build_chunk_pipeline()
        nodes = pipeline.run(documents=docs)

        # For every node, every row that contains "类别 NN" must also contain
        # "不少于 NN+1 年" — i.e. the row's key and value stay together.
        import re
        row_re = re.compile(r"类别\s*(\d{2}).*?不少于\s*(\d+)\s*年")
        for n in nodes:
            for line in n.text.split("\n"):
                m = row_re.search(line)
                if m:
                    cat = int(m.group(1))
                    yrs = int(m.group(2))
                    assert yrs == cat + 1, \
                        f"row misaligned in chunk: {line!r} (cat={cat}, yrs={yrs})"

        # And no chunk should *start* with a stray "| 不少于 N 年 |" row whose
        # key is missing — that's the exact pattern that caused the original bug.
        stray = re.compile(r"^\s*\|\s*不少于\s*\d+\s*年\s*\|\s*$")
        for n in nodes:
            first_nonempty = next((l for l in n.text.split("\n") if l.strip()), "")
            # "[breadcrumbs]" prepended by enricher is fine; check the body
            body = n.text.split("\n", 1)[1] if first_nonempty.startswith("[") else n.text
            for line in body.split("\n"):
                assert not stray.match(line), \
                    f"chunk starts with orphan trailing-cell row: {line!r}"

    def test_linked_list_neighbours(self, sample_md: Path):
        docs = parse_file(sample_md)
        for d in docs:
            d.metadata["doc_id"] = "doc-md-1"
        nodes = chunk_documents(docs)
        if len(nodes) < 2:
            pytest.skip("not enough chunks to test linkage")
        # First chunk: prev=None
        assert nodes[0].metadata["prev_chunk_id"] is None
        # Last chunk: next=None
        assert nodes[-1].metadata["next_chunk_id"] is None
        # Internal: prev/next form a chain
        for i in range(1, len(nodes) - 1):
            assert nodes[i].metadata["prev_chunk_id"] == nodes[i - 1].metadata["chunk_id"]
            assert nodes[i].metadata["next_chunk_id"] == nodes[i + 1].metadata["chunk_id"]
