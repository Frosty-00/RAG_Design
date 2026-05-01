"""PPTX reader — each slide is one section in the output Document."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.services.parsing.base import Document, make_base_metadata


class PptxReader:
    def load_data(self, file_path: str | Path) -> list[Document]:
        path = Path(file_path)
        prs = Presentation(str(path))
        chunks: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            title = None
            try:
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title = (slide.shapes.title.text_frame.text or "").strip()
            except Exception:  # noqa: BLE001
                title = None

            chunks.append(f"## Slide {i}" + (f" — {title}" if title else ""))

            body_lines: list[str] = []
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para in shape.text_frame.paragraphs:
                    text = "".join(r.text or "" for r in para.runs).strip()
                    if text:
                        body_lines.append(text)
            if body_lines:
                chunks.append("\n".join(body_lines))
            chunks.append("")  # blank between slides

        text = "\n".join(chunks).rstrip()
        if not text:
            return []
        return [Document(
            text=text,
            metadata=make_base_metadata(file_path=path, parser="pptx"),
        )]
